import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import io
import os
import html
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from streamlit_modal import Modal
from io import BytesIO
import base64
import pickle


# Function to load a PowerPoint presentation from BytesIO
def load_presentation_from_bytesio(presentation_bytesio):
    return Presentation(presentation_bytesio)

# Function to apply title styling
def style_title(title_shape):
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Default color black

# Function to apply title styling and background color
def set_slide_background_and_title_style(slide, title_shape, slide_index):
    
    # Create presentation object
    prs = Presentation()
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x06, 0x35, 0x7A)
    
    #Add slide number to title slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = 0  # Adjust the position 
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, slide_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x06, 0x35, 0x7A)  # Dark blue color
    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x06, 0x35, 0x7A) # Dark blue outline color
    
    text_frame = shape_blue.text_frame  
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right

    # Apply title styling
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color
            
            
# Function to clone shapes (including GraphicFrame and pictures)
def clone_shape(original_shape, merged_slide):
    if original_shape.shape_type == 14:  # Placeholder
        ph = original_shape.placeholder_format
        merged_slide.placeholders[ph.idx].text = original_shape.text
    elif original_shape.shape_type == 13:  # Picture
        image_stream = BytesIO(original_shape.image.blob)
        left = original_shape.left
        top = original_shape.top
        width = original_shape.width
        height = original_shape.height
        merged_slide.shapes.add_picture(image_stream, left, top, width, height)
    else:  # Other shapes (including GraphicFrame)
        merged_shape = original_shape._element
        merged_slide.shapes._spTree.append(merged_shape)

# Function to highlight Gini value based on user-defined thresholds
def highlight_gini(val, thresholds_gini):
    val = float(val)
    if val > thresholds_gini['green_gini']['value']:
        color = 'green'
    elif thresholds_gini['amber_gini']['lower'] < val <= thresholds_gini['amber_gini']['upper']:
        color = 'orange'
    elif val <= thresholds_gini['red_gini']['value']:
        color = 'red'
    else:
        color = 'white'
    return f'background-color: {color}'

# # Function to create an Excel file with highlighted cells
# def to_excel_with_highlights_gini(df, thresholds_gini):
#     output = io.BytesIO()
#     workbook = Workbook()
#     sheet = workbook.active

#     # Write the header
#     for col_num, column_title in enumerate(df.columns, 1):
#         cell = sheet.cell(row=1, column=col_num)
#         cell.value = column_title

#     # Write the data
#     for row_num, row in enumerate(df.values, 2):
#         for col_num, value in enumerate(row, 1):
#             cell = sheet.cell(row=row_num, column=col_num)
#             cell.value = value

#     # Highlight the last cell in the 'Gini Area' column
#     last_gini_cell = sheet.cell(row=len(df) + 1, column=df.columns.get_loc('Gini Area') + 1)
#     gini_value = df.iloc[-1]['Gini Area']
#     if gini_value > thresholds_gini['green_gini']['value']:
#         last_gini_cell.fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")
#     elif thresholds_gini['amber_gini']['lower'] < gini_value <= thresholds_gini['amber_gini']['upper']:
#         last_gini_cell.fill = PatternFill(start_color="FFBF00", end_color="FFBF00", fill_type="solid")
#     elif gini_value <= thresholds_gini['red_gini']['value']:
#         last_gini_cell.fill = PatternFill(start_color="FF0000", end_color="FF0000", fill_type="solid")

#     workbook.save(output)
#     return output.getvalue()

# Function to add a slide with the ribbon and logo
def ppt_ribbon_and_logo(slide, slide_index):
    # Create presentation object
    prs = Presentation()
    # Determine slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate ribbon dimensions
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2

    # Calculate positions for ribbons
    ribbon_left = Inches(0)
    ribbon_top = slide_height - ribbon_height

    # Add the first half of the ribbon (amber)
    shape_amber = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_amber = shape_amber.fill
    fill_amber.solid()
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        fill_amber.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber color
    else:
        fill_amber.fore_color.rgb = RGBColor(255, 191, 0)  # Amber color

    line_amber = shape_amber.line
    # Check if color is in session state; if not, use default color
    if 'ribbon_color_1' in st.session_state:
        line_amber.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_1[1:])  # Amber outline color
    else:
        line_amber.color.rgb = RGBColor(255, 191, 0)  # Amber color

    # Add text to the amber ribbon
    text_frame_amber = shape_amber.text_frame
    text_frame_amber.text = "\tENBD Model Monitoring"
    text_frame_amber.paragraphs[0].font.size = Pt(9)
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_font_1' in st.session_state:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_1[1:])  # Black text
    else:
        text_frame_amber.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

    # Add the second half of the ribbon (dark blue)
    ribbon_left = half_ribbon_width  # Adjust the left position for the second half
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, half_ribbon_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if olor is in session state; if not, use default color
    if 'ribbon_color_2' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x06, 0x35, 0x7A)  # Dark blue color

    line_blue = shape_blue.line
    # Check if title_font_color is in session state; if not, use default color black
    if 'ribbon_color_2' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.ribbon_color_2[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x06, 0x35, 0x7A) # Dark blue outline color
    
    # Add slide number
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    half_ribbon_width = slide_width / 2
    ribbon_top = slide_height - ribbon_height
    
    text_frame = shape_blue.text_frame
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_2' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_2[1:])  # White text
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
    

    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD_s.jpg')  # Replace with your actual file path
    logo_left = slide_width - Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.15)  # Adjust position as needed
    logo_height = Inches(0.45)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    
#Creating ppt for gini
def create_ppt_gini(df, fig_bytes, thresholds_gini, data_comment, graph_comment):
    prs = Presentation()
    
    slide_index = 1  # To keep track of the slide index

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Set the background color to the slide
    background = slide.background
    fill = background.fill
    fill.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])
    else:
        fill.fore_color.rgb = RGBColor(0x06, 0x35, 0x7A)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "PL - Scorecard Model Gini"
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.6)  # Centered vertically
    title_shape.width = Inches(8)
    title_shape.height = Inches(2)
    
    # Add logo image to the upper right corner
    base_dir = os.path.dirname(__file__)
    
    # Logo path
    logo_path = os.path.join(base_dir, 'Images', 'ENBD.jpg')  # Replace with your actual file path
    logo_left = Inches(0.6)  # Adjust position as needed
    logo_top = Inches(0.25)  # Adjust position as needed
    logo_height = Inches(0.6)  # Adjust size as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(42)
            run.font.bold = True
            # Check if color is in session state; if not, use default color
            if 'font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.font_color[1:])  # White font color
            else:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color
         
   # Add slide number to title slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ribbon_height = Inches(0.28)  # Adjust height as needed
    ribbon_top = slide_height - ribbon_height
    
    # Add the second half of the ribbon (dark blue)
    ribbon_left = 0  # Adjust the position
    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ribbon_left, ribbon_top, slide_width, ribbon_height
    )
    fill_blue = shape_blue.fill
    fill_blue.solid()
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        fill_blue.fore_color.rgb = RGBColor.from_string(st.session_state.bg_color[1:])  # Dark blue color
    else:
        fill_blue.fore_color.rgb = RGBColor(0x06, 0x35, 0x7A)  # Dark blue color
    line_blue = shape_blue.line
    # Check if color is in session state; if not, use default color
    if 'bg_color' in st.session_state:
        line_blue.color.rgb = RGBColor.from_string(st.session_state.bg_color[1:]) # Dark blue outline color
    else:
        line_blue.color.rgb = RGBColor(0x06, 0x35, 0x7A) # Dark blue outline color
    
    text_frame = shape_blue.text_frame    
    text_frame.text = f" {slide_index}\t"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(9)
    # Check if color is in session state; if not, use default color
    if 'ribbon_font_color' in st.session_state:
        p.font.color.rgb = RGBColor.from_string(st.session_state.ribbon_font_color[1:])
    else:
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
    p.alignment = PP_ALIGN.RIGHT  # Align text to the right
            
    slide_index += 1

    # Data Table slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    
    ppt_ribbon_and_logo(slide, slide_index)

    # Add and style title with a shorter height
    title_shape = shapes.title
    title_shape.text = "Gini calculation"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Determine table size and position
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)  # Adjusted to fit the table within the slide

    rows, cols = df.shape
    table = shapes.add_table(rows + 1, cols, left, top, width, height).table
    table_style_id = table._tbl.tblPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
    )
    table_style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"

    # Set column names and font size
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        # Check if color is in session state; if not, use default color
        if 'row_bg_color' in st.session_state:
            cell.fill.fore_color.rgb = RGBColor.from_string(st.session_state.row_bg_color[1:]) # Teal
        else:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80) # Teal
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                # Check if color is in session state; if not, use default color
                if 'row_font_color' in st.session_state:
                    run.font.color.rgb = RGBColor.from_string(st.session_state.row_font_color[1:])  # White font color
                else:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                
    # Add data to table and set font size
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            # Format decimal values to 3 decimal places
            if isinstance(value, float):
                cell.text = f"{value:.4f}".rstrip('0').rstrip('.')
            else:
                cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    # Check if color is in session state; if not, use default color
                    if 'content_font_color' in st.session_state:
                        run.font.color.rgb = RGBColor.from_string(st.session_state.content_font_color[1:])
                    else:
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        
            if row_idx == len(df) - 1 and col_idx == df.columns.get_loc('Gini Area'):
                if value > thresholds_gini['green_gini']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
                elif thresholds_gini['amber_gini']['lower'] < value <= thresholds_gini['amber_gini']['upper']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xBF, 0x00)
                elif value <= thresholds_gini['red_gini']['value']:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # Set row heights
    for row in range(rows + 1):
        table.rows[row].height = Inches(0.3)

    # Add data comment
    data_comment_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1))
    text_frame = data_comment_box.text_frame
    text_frame.text = f"Comment: {data_comment}"
    # Set font size to 10 for all paragraphs in the text frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
    
    slide_index += 1

    # Chart slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    
    ppt_ribbon_and_logo(slide, slide_index)
    
    title_shape = shapes.title
    title_shape.text = "Graph"
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9)
    title_shape.height = Inches(0.5)
    title_text_frame = title_shape.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(22)
            # Check if color is in session state; if not, use default color
            if 'title_font_color' in st.session_state:
                run.font.color.rgb = RGBColor.from_string(st.session_state.title_font_color[1:])
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Add chart image
    image_stream = io.BytesIO(fig_bytes)
    shapes.add_picture(image_stream, Inches(0.6), Inches(0.8), Inches(8.8), Inches(4.5))

    # Add graph comment
    graph_comment_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    text_frame = graph_comment_box.text_frame
    text_frame.text = f"Comment: {graph_comment}"
    # Set font size to 10 for all paragraphs in the text frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    # Save presentation
    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    return ppt_output


# Function to create the threshold selection UI for Gini and save thresholds
def threshold_selection_gini(show_ui=True):
    # Base directory of the current script
    base_dir = os.path.dirname(__file__)

    # Construct the path for the pickle file
    file_path = os.path.join(base_dir, 'pkl', 'model_gini.pkl')

    try:
        with open(file_path, 'rb') as f:
            thresholds_gini = pickle.load(f)
    except FileNotFoundError:
        thresholds_gini = None

    if thresholds_gini:
        green_threshold = thresholds_gini['green_gini']['value']
        amber_lower = thresholds_gini['amber_gini']['lower']
        amber_upper = thresholds_gini['amber_gini']['upper']
        red_threshold = thresholds_gini['red_gini']['value']
    else:
        # Initialize with None if no thresholds are found
        green_threshold = None
        amber_lower = None
        amber_upper = None
        red_threshold = None

    if show_ui:
        with st.expander('Please select the threshold values for Gini'):
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)
            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<p style="font-size:17px;"><b>Legend</b></p>', unsafe_allow_html=True)
            with c2:
                st.markdown('<p style="font-size:17px;"><b>Thresholds</b></p>', unsafe_allow_html=True)
            with c3:
                st.markdown('<p style="font-size:17px;"><b>Evaluation</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)
    
            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<div style="background-color: green; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                green_threshold = st.number_input("Green Threshold (Gini Area > 40%)", value=green_threshold or 0.40)
            with c3:
                st.markdown('<p style="font-size:19px;"><b>No Action Required</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)
    
            c1, c2, c3, c4 = st.columns(4)
            with c1:
                st.markdown('<div style="background-color: orange; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                amber_lower = st.number_input("Amber Lower Limit (30% < Gini Area <= 40%)", value=amber_lower or 0.30)
            with c3:
                amber_upper = st.number_input("Amber Upper Limit (30% < Gini Area <= 40%)", value=amber_upper or 0.40)
            with c4:
                st.markdown('<p style="font-size:19px;"><b>To be Discussed</b></p>', unsafe_allow_html=True)
            st.markdown('<hr style="margin-top: -5px; margin-bottom: -5px;">', unsafe_allow_html=True)
    
            c1, c2, c3 = st.columns([2, 4, 2])
            with c1:
                st.markdown('<div style="background-color: red; width: 80%; height: 40px;"></div>', unsafe_allow_html=True)
            with c2:
                red_threshold = st.number_input("Red Threshold (Gini Area <= 30%)", value=red_threshold or 0.30)
            with c3:
                st.markdown('<p style="font-size:19px;"><b>Action Required</b></p>', unsafe_allow_html=True)

        # Save thresholds to file
        threshold = {
            'green_gini': {'value': green_threshold},
            'amber_gini': {'lower': amber_lower, 'upper': amber_upper},
            'red_gini': {'value': red_threshold}
        }
        # Base directory of the current script
        base_dir = os.path.dirname(__file__)
        
        # Construct the path for the pickle file
        file_path = os.path.join(base_dir, 'pkl', 'model_gini.pkl')
        
        with open(file_path, 'wb') as f:
            pickle.dump(threshold, f)

    return {
        'green_gini': {'value': green_threshold},
        'amber_gini': {'lower': amber_lower, 'upper': amber_upper},
        'red_gini': {'value': red_threshold}
    }

def create_ppt_download_button_gini(df, fig_bytes, thresholds_gini, data_comment="", graph_comment=""):
    
    # Create PowerPoint presentation bytes
    ppt_data_gini = create_ppt_gini(df, fig_bytes, thresholds_gini, data_comment, graph_comment)
    
    return ppt_data_gini

# Function to merge four PowerPoint presentations and apply title styling to specific slides
def merge_presentations(presentation1, presentation2, presentation3, presentation4):
    merged_presentation = Presentation()
    slide_index = 1  # To keep track of the slide index

    # Copy slides from presentation1 (Overview) first
    for slide in presentation1.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [1]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation2 (Change Log)
    for slide in presentation2.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 1
        if merged_slide.shapes.title:
            if slide_index in [2]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
        slide_index += 1
        
    # Copy slides from presentation3 (Summary)
    for slide in presentation3.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 3
        if merged_slide.shapes.title:
            if slide_index in [3]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
      
        slide_index += 1
    
    # Copy slides from presentation4 (Gini)
    for slide in presentation4.slides:
        slide_layout = merged_presentation.slide_master.slide_layouts[5]
        merged_slide = merged_presentation.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            clone_shape(shape, merged_slide)

        # Apply title styling if there is a title and the slide index is 5 or 6
        if merged_slide.shapes.title:
            if slide_index in [5, 6]:
                style_title(merged_slide.shapes.title)
                ppt_ribbon_and_logo(merged_slide, slide_index)
        
            if slide_index in [4]:
                set_slide_background_and_title_style(merged_slide, merged_slide.shapes.title, slide_index)
    
        slide_index += 1
    
    return merged_presentation

# Streamlit app
def app():
    st.markdown(
                """
                <h1 style='text-align: center; font-size: 28px; color: rgb(39, 45, 85);'>PL - Scorecard Model Gini</h1>
                """,
                unsafe_allow_html=True
            )

    # Base directory of the current script
    base_dir = os.path.dirname(__file__)
    
    # Read the Excel file
    df = pd.read_excel(os.path.join(base_dir, 'Datasets', 'Gini_Data_dashboard.xlsx'))   # Ensure the file 'data.xlsx' is in the same directory as your script
    # Replace None values with empty strings for better display
    df1 = df.fillna("")
    st.session_state.df_gini = df  # Save df to session_state
    
    
    # Initialize comments
    data_comment_gini = ""
    graph_comment_gini = ""
    
    # Create two tabs: one for data and one for the graph
    tab1, tab2 = st.tabs(["Gini Calculation", "Graph"])
    
    with tab1:
        thresholds_gini = threshold_selection_gini(show_ui=True)
        st.session_state.thresholds_gini = thresholds_gini
        
        st.markdown(
            """
            <div style='text-align: center;
            font-size: 20px;'>
                <strong>Gini Result</strong>
            </div>
            """,
            unsafe_allow_html=True)
        
        
        # def create_excel_download_button_gini(df, thresholds_gini):
        #     # Convert DataFrame to Excel with highlights
        #     excel_data = to_excel_with_highlights_gini(df, thresholds_gini)
            
        #     # Save byte data into st.session_state
        #     st.session_state.gini_workbook_data = excel_data
            
        #     # Create a buffer for the Excel data
        #     buffer = BytesIO()
        #     buffer.write(excel_data)
        #     buffer.seek(0)
            
        #     # Encode the file in base64
        #     b64 = base64.b64encode(buffer.read()).decode()
            
        #     # Read and encode the Excel image to base64
        #     with open("excel_logo.png", "rb") as image_file:
        #         excel_image_base64 = base64.b64encode(image_file.read()).decode()
            
        #     # Define the HTML for the download icon link
        #     download_html = f'''
        #         <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}" 
        #            download="highlighted_data.xlsx" class="excel-download-button" title="Click here to download the Excel file">
        #            <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
        #         </a>
        #     '''
            
        #     return download_html
        
        # download_html = create_excel_download_button_gini(df, thresholds_gini)
       
        # # Display the download link in Streamlit
        # st.markdown(download_html, unsafe_allow_html=True)
        
    #     # Apply styling to the DataFrame
    #     styled_df = df.style.applymap(lambda x: highlight_gini(x, thresholds_gini) if x == df['Gini Area'].iloc[-1] else '', subset=['Gini Area'])

    #    # Display the DataFrame without unnecessary trailing zeros
    #     df_styled =styled_df.format(lambda x: f"{x:.4f}".rstrip('0').rstrip('.') if isinstance(x, float) else f"{x}")

        # Function for formatting values
        def format_value(val):
            if isinstance(val, float):
                # Format floats to 4 decimal places, remove trailing zeros and the dot if not needed
                return f"{val:.4f}".rstrip('0').rstrip('.')
            return str(val)

        # Function for creating HTML table with conditional styling
        def create_html_table_with_download(dataframe, file_name, image_path):
            # Encode the image as base64
            with open(image_path, "rb") as image_file:
                image_base64 = base64.b64encode(image_file.read()).decode()

            # Create download link
            download_link = f"""
            <style>
                .download-icon img {{
                    width: 70px;
                    height: 70px;
                    border-radius: 30%;  /* Makes the image circular */
                    transition: transform 0.3s ease, box-shadow 0.3s ease;
                }}
                .download-icon img:hover {{
                    transform: scale(1.1);  /* Slight zoom effect on hover */
                    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.3);  /* Shadow effect */
                }}
                .download-icon img:active {{
                    transform: scale(0.90);  /* Slight shrink effect on click */
                    box-shadow: none;  /* Remove shadow when clicked */
                }}
            </style>

            <div class="download-icon">
                <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(to_excel(dataframe)).decode()}" download="{file_name}" title="Click to download the file">
                    <img src="data:image/png;base64,{image_base64}" alt="Download Icon" style="width:27px; height:auto;">
                </a>
            </div>
            """

            last_row_index = dataframe.index[-1]

            # Create HTML table with inline styling
            html_table = f"""
            <div class="custom-container">
                <table class="dataframe">
                    <thead><tr>
            """
            for col_name in dataframe.columns:
                html_table += f'<th>{html.escape(col_name)}</th>'
            html_table += '</tr></thead><tbody>'
            for index, row in dataframe.iterrows():
                html_table += '<tr>'
                for col_name, col_value in row.items():
                    style = ''
                    formatted_value = format_value(col_value)  # Apply formatting
                    if col_name == 'Gini Area' and index == last_row_index:
                        style = highlight_gini(col_value, thresholds_gini)
                    html_table += f'<td style="{style}">{html.escape(formatted_value)}</td>'
                html_table += '</tr>'
            html_table += '</tbody></table></div>'

            # Combine the download link and the table
            return download_link + html_table

        # Function to convert DataFrame to Excel
        def to_excel(df):
            output = BytesIO()
            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Sheet1')
            return output.getvalue()

        # Optional: If you want to display custom CSS separately
        custom_css = """
        <style>
            .custom-container {
                max-height: 400px;
                max-width: 100%;
                overflow-y: scroll;
                overflow-x: scroll;
                position: relative;
                border-radius: 10px;
                border: 1px solid #ccc;
            }
            table {
                width: 100%;
                height: auto;
            }
            th, td {
                font-size: 14px;
                padding: 8px;
                text-align: left;
                white-space: nowrap;  /* Prevent text from wrapping */

            }
            .download-icon {
                position: absolute;
                right: 40px;
                top: -40px;  /* Adjust to align vertically outside the table container */
                font-size: 24px;   
            }
            thead {
                position: sticky;
                top: 0;
                background-color: rgb(39, 45, 85);
                color: rgb(20, 26, 63);
            }
            thead th {
                color: white;
                font-weight: normal;
            }
            .info-container {
                margin-top: 10px;
                padding: 8px;
                border: 1px solid #ccc;
                border-radius: 5px;
                background-color: #f9f9f9;
                font-size: 14px;
                text-align: left;
                color: black; 
            }
        </style>
        """
        st.markdown(custom_css, unsafe_allow_html=True)

        # Base directory of the current script
        base_dir = os.path.dirname(__file__)

        # Path to the image file
        image_path = os.path.join(base_dir, 'Images', 'Download_icon.png')

        # Create HTML table and display it
        html_table = create_html_table_with_download(df, "Gini_result.xlsx", image_path)
        st.markdown(html_table, unsafe_allow_html=True)
        
        # Display the DataFrame with highlighted cell
        # st.dataframe(df_styled, width=1200)
        
        # Add space between the table and download button
        st.markdown(
            """
            <style>
                .spacer {
                    margin-top: 20px; /* Adjust the value to increase or decrease the space */
                }
            </style>
            """,
            unsafe_allow_html=True,
        )
        
        # Add comment box for data
        data_comment_modal = Modal("Comment", key="data_comment")
        if st.button("Add Comment", key="data_comment_button", help = "Click here to add Comment"):
            data_comment_modal.open()
            
        if data_comment_modal.is_open():
            with data_comment_modal.container():
                data_comment_gini = st.text_area("Enter your comment:", key="data_comment_textarea")
                if st.button("Submit Comment", key="submit_data_comment"):
                    st.session_state.data_comment_gini = data_comment_gini
                    data_comment_modal.close()
        
# =============================================================================
#         # Prepare Excel file in memory
#         excel_data = io.BytesIO()
#         with pd.ExcelWriter(excel_data, engine='openpyxl') as writer:
#             df_styled.to_excel(writer, index=False, sheet_name='Gini', startrow=12)
#         excel_data.seek(0)
#         
#         # Save byte data into st.session_state
#         st.session_state.gini_workbook_data = excel_data
#         
#         # Convert BytesIO to bytes
#         excel_data_bytes = excel_data.getvalue()
#         
#         # Encode the Excel data to base64
#         excel_data_base64 = base64.b64encode(excel_data_bytes).decode()
#         
#         # Read and encode the Excel image to base64
#         with open("excel_logo.png", "rb") as image_file:
#             excel_image_base64 = base64.b64encode(image_file.read()).decode()
#         
#         # Define the HTML for the download icon link
#         download_html = f'''
#             <a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_data_base64}" 
#                download="Gini.xlsx" class="excel-download-button" title="Click here to download the Excel file">
#                <img src="data:image/png;base64,{excel_image_base64}" alt="Download Excel">
#             </a>
#         '''
#         
#         # Display the custom Excel download button in Streamlit
#         st.markdown(download_html, unsafe_allow_html=True)
# =============================================================================
        custom_css = """
            <style>
                .excel-download-button {
                    position: absolute;
                    top: -72px;
                    left: 150px;
                    cursor: pointer;
                    padding: 5.5px 12px;  /* Adjust padding as needed (top/bottom, left/right) */
                }
            </style>
        """
        st.markdown(custom_css, unsafe_allow_html=True)

        # Ensure that you check if the button HTML is available in session state
        if 'excel_button_html' in st.session_state:
            st.markdown(st.session_state.excel_button_html, unsafe_allow_html=True)
        else:
            st.write("No data available for download in Excel. Please run the Overview and Change Log modules first.")
        
        
    with tab2:
        st.markdown(
                """
                <div style='text-align: center;
                font-size: 20px;'>
                    <strong>Graph</strong>
                </div>
                """,
                unsafe_allow_html=True)
                    
        # Assuming df and df1 are already defined dataframes
        fig = go.Figure()
        
        # Add traces
        fig.add_trace(go.Scatter(x=df['% Cum Total'], y=df['% Cum Total'], mode='lines', name='Random', line=dict(color='red', width = 2)))
        fig.add_trace(go.Scatter(x=df['% Cum Total'], y=df['% Cum Bad'], mode='lines', name='Actual', line=dict(color='blue', width = 2)))
        fig.add_trace(go.Scatter(x=df['Perfect Curve'], y=df['Perfect Curve_1'], mode='lines', name='Perfect Curve', line=dict(color='purple', width = 2)))
        fig.add_trace(go.Scatter(x=df1['% Cum Total'], y=df1['% Cum Bad'], mode='lines', name='Development', line=dict(color='gray', width = 2)))
        
        # Update layout
        fig.update_layout(
            title='',
            xaxis_title=dict(
                text='Cumulative Total -->',
                font=dict(
                    size=17,
                    color='black',
                    family='Calibri'
                )
            ),
            yaxis_title=dict(
                text='Cumulative Bad -->',
                font=dict(
                    size=17,
                    color='black',
                    family='Calibri'
                )
            ),
            xaxis=dict(
                showgrid=True,  # Show gridlines
                gridcolor='lightgray',  # Color of gridlines
                showline=True, 
                linecolor='black', 
                linewidth=2, mirror=True
            ),
            yaxis=dict(
                showgrid=True,  # Show gridlines
                gridcolor='lightgray',  # Color of gridlines
                showline=True, 
                linecolor='black', 
                linewidth=2, mirror=True
            ),
            legend=dict(
                orientation='h', x=0.3, y=-0.2,
                bordercolor='black', borderwidth=1
            ),
            margin=dict(l=60, r=30, t=40, b=60),
            showlegend=True,
            width = 1200, 
            height = 500
        )
        
        st.plotly_chart(fig, use_container_width=True)
        
        # Add comment box for graph
        graph_comment_modal = Modal("Comment", key="graph_comment")
        if st.button("Add Comment", key="graph_comment_button", help = "Click here to add Comment"):
            graph_comment_modal.open()
            
        if graph_comment_modal.is_open():
            with graph_comment_modal.container():
                graph_comment_gini = st.text_area("Enter your comment:", key="graph_comment_textarea")
                if st.button("Submit Comment", key="submit_graph_comment"):
                    st.session_state.graph_comment_gini = graph_comment_gini
                    graph_comment_modal.close()
                
        # Convert plot figure to PNG bytes
        fig_bytes = pio.to_image(fig, format='png')
        st.session_state.fig_bytes = fig_bytes
        st.session_state.thresholds_gini = thresholds_gini
        
        ppt_data_gini = create_ppt_download_button_gini(df, fig_bytes, thresholds_gini, st.session_state.get("data_comment_gini", ""), st.session_state.get("graph_comment_gini", ""))
        
        ppt_data_overview = ppt_data_change_log = ppt_data_summary = None
        
        if 'ppt_data_overview' in st.session_state:
            ppt_data_overview = st.session_state.ppt_data_overview
        
        if 'ppt_data_change_log' in st.session_state:
            ppt_data_change_log = st.session_state.ppt_data_change_log
            
        if 'ppt_data_summary' in st.session_state:
            ppt_data_summary = st.session_state.ppt_data_summary
            
        if ppt_data_overview and ppt_data_change_log and ppt_data_summary and ppt_data_gini:
            presentation1 = load_presentation_from_bytesio(ppt_data_overview)
            presentation2 = load_presentation_from_bytesio(ppt_data_change_log)
            presentation3 = load_presentation_from_bytesio(ppt_data_summary)
            presentation4 = load_presentation_from_bytesio(ppt_data_gini)
            
            merged_presentation = merge_presentations(presentation1, presentation2, presentation3, presentation4)
            merged_presentation_bytesio = BytesIO()
            merged_presentation.save(merged_presentation_bytesio)
            merged_presentation_bytesio.seek(0)
        
            # Base directory of the current script
            base_dir = os.path.dirname(__file__)
            
            # Construct the path for the image
            image_path = os.path.join(base_dir, "Images", "ppt_logo.png")
            
            # Read and encode the image to base64
            with open(image_path, "rb") as image_file:
                image_base64 = base64.b64encode(image_file.read()).decode()
            
            # Convert BytesIO to bytes
            ppt_data_merged_bytes = merged_presentation_bytesio.getvalue()
            
            # Encode the PPT data to base64
            ppt_data_merged_base64 = base64.b64encode(ppt_data_merged_bytes).decode()
            
            # Help text
            help_text = "Click here to download the Dashboard into PowerPoint presentation"
            
            custom_css = """
            <style>
                .ppt-download-button {
                    position: absolute;
                    top: -72px;
                    left: 150px;
                    cursor: pointer
                }
            </style>
            """
            st.markdown(custom_css, unsafe_allow_html=True)
            
            # Create the HTML for the button with image and help text
            ppt_button_html = f"""
            <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_data_merged_base64}" download="Dashboard.pptx" class="ppt-download-button" title="{help_text}">
                <img src="data:image/png;base64,{image_base64}" alt="Download PPT">
            </a>
            """
            
            # Display the custom download button in Streamlit
            st.markdown(ppt_button_html, unsafe_allow_html=True)
            


if __name__ == "__main__":
    app()
